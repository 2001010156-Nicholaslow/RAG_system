import os
import uuid
import pymupdf
import hashlib
import docx
import json
import pandas as pd
import xml.etree.ElementTree as ET
from dotenv import load_dotenv
from bs4 import BeautifulSoup
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from psycopg2 import pool, sql
from pathlib import Path
from typing import List, Optional, Callable, Dict, Any
from functools import lru_cache
from tqdm import tqdm
from datetime import datetime
import logging
import argparse
from transformers import AutoTokenizer

# === CONFIG ===
EMBEDDING_MODEL_NAME = "Qwen/Qwen3-Embedding-0.6B"
TOKENIZER_NAME = "Qwen/Qwen3-Embedding-0.6B"

logging.basicConfig(
    level=logging.INFO,
    format='[%(asctime)s] %(levelname)s - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)

logger = logging.getLogger(__name__)

load_dotenv()

# === Database Connection Pool ===
db_pool = pool.SimpleConnectionPool(
    1, 10,
    database=os.getenv("DB_NAME"),
    user=os.getenv("DB_USER"),
    password=os.getenv("DB_PASSWORD"),
    host=os.getenv("DB_HOST"),
    port=os.getenv("DB_PORT"),
)

def get_db_connection():
    return db_pool.getconn()

def release_db_connection(conn):
    db_pool.putconn(conn)

# === Lazy Load Embedding Model ===
_model = None
tokenizer = AutoTokenizer.from_pretrained(TOKENIZER_NAME)

def get_model():
    global _model
    if _model is None:
        _model = SentenceTransformer(EMBEDDING_MODEL_NAME)
    return _model

# === Chunking Function ===
def chunk_text(
        text: str,
        tokenizer: Optional[Any] = None, 
        max_tokens: int = 500,
        overlap_tokens: int = 50,
        min_chunk_size: int = 50
        ) -> List[str]:
    paragraphs = [p.strip() for p in text.splitlines() if p.strip()]
    chunks, current_chunk, current_tokens = [], [], 0

    @lru_cache(maxsize=1000)
    def count_tokens(p: str) -> int:
        return len(tokenizer.tokenize(p)) if tokenizer else len(p.split())

    def split_large_para(para: str) -> List[str]:
        if tokenizer:
            tokens = tokenizer.tokenize(para)
            return [
                tokenizer.convert_tokens_to_string(tokens[i:i+max_tokens]) 
                for i in range(0, len(tokens), max_tokens - overlap_tokens)]
        return [
            para[i:i+max_tokens]
            for i in range(0, len(para), max_tokens - overlap_tokens)]

    for para in paragraphs:
        para_tokens = count_tokens(para)
        if para_tokens > max_tokens:
            if current_chunk:
                chunks.append('\n'.join(current_chunk))
                current_chunk, current_tokens = [], 0
            chunks.extend(split_large_para(para))
            continue

        if current_tokens + para_tokens > max_tokens:
            if current_chunk:
                chunks.append('\n'.join(current_chunk))
                overlap, overlap_size = [], 0
                for p in reversed(current_chunk):
                    p_tokens = count_tokens(p)
                    if overlap_size + p_tokens <= overlap_tokens:
                        overlap.insert(0, p)
                        overlap_size += p_tokens
                    else:
                        break
                current_chunk = overlap
                current_tokens = overlap_size
            else:
                current_chunk = []
                current_tokens = 0

        current_chunk.append(para)
        current_tokens += para_tokens

    if current_chunk and current_tokens >= min_chunk_size:
        chunks.append('\n'.join(current_chunk))
    return chunks

# === Embedding Chunks ===
def embed_chunks(chunks: List[str]) -> List[List[float]]:
    model = get_model()
    logger.info("Embedding chunks...")
    return model.encode(chunks, normalize_embeddings=True, batch_size=16).tolist()

# === Metadata Generator ===
def generate_metadata(
        file_path: str,
        chunk_count: int,
        document_name :str,
        permissions: Optional[List[str]],
        file_type: Optional[str],
        category: Optional[str],
        ) -> Dict[str, Any]:
    stat = os.stat(file_path)
    with open(file_path, "rb") as f:
        file_hash = hashlib.md5(f.read()).hexdigest()
    return {
        "file_size_bytes": stat.st_size,
        "file_type": file_type,
        "chunk_count": chunk_count,
        "file_hash": file_hash,
        "date_updated": datetime.now().isoformat(),
        "permission": permissions or [],
        "category": category or "Uncategorized",
        "embedding_model": EMBEDDING_MODEL_NAME,
        "tokenizer" : TOKENIZER_NAME,
        "document_name": document_name
    }

# === Insert into DB ===
def insert_chunks(
        chunks: List[str],
        document_name: str,
        file_path: str,
        permissions: Optional[List[str]],
        category: Optional[str]
        ):
    conn = None
    try:
        logger.info("Starting Database coneection...")
        conn = get_db_connection()
        cur = conn.cursor()
        permissions = permissions
        embeddings = embed_chunks(chunks)
        file_type = os.path.splitext(file_path)[1].lstrip(".")
        metadata = generate_metadata(file_path, len(chunks), document_name, permissions, file_type, category)
        table_name = "documents"
        file_hash = metadata["file_hash"]

        cur.execute(sql.SQL("SELECT COUNT(*) FROM {} WHERE metadata ->> 'file_hash' = %s")
                    .format(sql.Identifier(table_name)),
                    (file_hash,))
        if cur.fetchone()[0] > 0:
            logger.warning(f"Duplicate file detected (hash: {file_hash}). Skipping insert.")
            return

        data = []
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
            data.append((
                str(uuid.uuid4()),
                chunk,
                embedding,
                document_name,
                permissions,
                i,
                json.dumps(metadata),
                file_hash
            ))
        cur.executemany(sql.SQL("""
            INSERT INTO {} (
                id, chunk, embedding, document_name,
                permissions, chunk_index, metadata, file_hash
            ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
        """).format(sql.Identifier(table_name)), data)

        conn.commit()
        logger.info(f"✅ Inserted {len(data)} chunks into {table_name}.")
    except Exception as e:
        logger.error(f"❌ DB error: {e}")
        if conn:
            conn.rollback()
        raise
    finally:
        if conn:
            release_db_connection(conn)

# === File Reading ===
def read_pdf(path) -> Optional[str]:
    try:
        with pymupdf.open(path) as doc:
            text = "\n".join([page.get_text() for page in doc]) # type: ignore
        return text
    except Exception as e:
        logger.error(f"Error reading PDF: {e}")
        return None

def read_docx(path) -> Optional[str]:
    try:
        doc = docx.Document(path)
        texts = [p.text for p in doc.paragraphs]
        # Add table content
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts.extend(p.text for p in cell.paragraphs)
        
        text = "\n".join(filter(None, texts))  # Remove empty strings
        return text if text.strip() else None
    except Exception as e:
        logger.error(f"Error reading docx: {e}")
        return None

def read_txt(path) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()
            return text
    except UnicodeDecodeError:
        try:
            with open(path, "r", encoding='latin-1') as f:
                text = f.read()
                return text if text.strip() else None
        except Exception as e:
            logger.error(f"Error reading txt: {e}")
            return None
    except Exception as e:
        logger.error(f"Error reading txt: {e}")
        return None

def read_excel(path) -> Optional[str]:
    try:
        dfs = pd.read_excel(path, sheet_name=None)
        return "\n\n".join(
            f"Sheet: {name}\n{df.to_string(index=False)}" 
            for name, df in dfs.items()
        )
    except Exception as e:
        logger.error(f"Error reading excel: {e}")
        return None

def read_json(path) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Error reading JSON: {e}")
        return None

def read_csv(path) -> Optional[str]:
    try:
        df = pd.read_csv(path)
        return df.to_string(index=False)
    except Exception as e:
        logger.error(f"Error reading CSV: {e}")
        return None

def read_xml(path) -> Optional[str]:
    try:
        tree = ET.parse(path)
        ET.indent(tree)
        return ET.tostring(tree.getroot(), encoding="unicode")
    except Exception as e:
        logger.error(f"Error reading XML: {e}")
        return None

def read_html(path: str) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        
        # Remove script and style elements
        for element in soup(["script", "style"]):
            element.decompose()
            
        text = soup.get_text(separator="\n", strip=True)
        return text if text.strip() else None
    except Exception as e:
        logger.error(f"Error reading HTML: {e}")
        return None


def read_pptx(path) -> Optional[str]:
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs: # type: ignore
                    for run in paragraph.runs:
                        texts.append(run.text)
        
        if not texts:  # No text found in presentation
            return None
        
        pptx_result = "\n".join(texts)
        return pptx_result
    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        return None

# === Unified File Reader ===
def read_file(path) -> Optional[str]:
    path = str(Path(path).resolve())
    if not os.path.exists(path):
        raise FileNotFoundError(f"File not found: {path}")

    ext = Path(path).suffix.lower().lstrip('.')
    reader_dispatch: dict[str, Callable[[str], Optional[str]]] = {
        "pdf": read_pdf,
        "docx": read_docx,
        "txt": read_txt,
        "xls": read_excel,
        "xlsx": read_excel,
        "json": read_json,
        "csv": read_csv,
        "xml": read_xml,
        "html": read_html,
        "pptx": read_pptx
    }
    reader = reader_dispatch.get(ext)
    if not reader:
        raise ValueError(f"Unsupported file type: {ext}")
    return reader(path)

# === CLI Entry Point ===
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True)
    parser.add_argument("--permissions", nargs="+", default=["default"])
    parser.add_argument("--category", default="Uncategorized")
    args = parser.parse_args()

    try:
        logger.info(f"Reading: {args.file}")
        full_text = read_file(args.file)
        if not full_text:
            raise ValueError("No text extracted from file.")

        logger.info("Chunking...")
        chunks = chunk_text(full_text)
        logger.info(f"{len(chunks)} chunks created. Inserting into DB...")

        insert_chunks(
            chunks=chunks,
            document_name=os.path.basename(args.file),
            file_path=args.file,
            permissions=args.permissions,
            category=args.category
        )
        logger.info("✅ Process Completed.")
    except Exception as e:
        logger.error(f"❌ Pipeline failed: {e}")
